"""
Comprehensive Report Generator for Weather Regulation Prediction

This module generates detailed analysis reports in multiple formats including:
- HTML reports with interactive visualizations
- PDF reports with static charts and tables
- LaTeX reports for academic papers
- Markdown reports for documentation
- PowerPoint presentations for stakeholders
"""

import os
import json
import pandas as pd
import numpy as np
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Union, Any, Tuple
import matplotlib.pyplot as plt
import seaborn as sns
from jinja2 import Template, Environment, FileSystemLoader
import base64
from io import BytesIO
import warnings
warnings.filterwarnings('ignore')

# Optional imports for PDF and PowerPoint generation
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from results.results_manager import ExperimentResult, ModelResult
from visualization.plots import ModelVisualizer, WeatherVisualizer


class ReportGenerator:
    """Generate comprehensive analysis reports"""
    
    def __init__(self, template_dir: Optional[str] = None):
        """Initialize report generator"""
        self.template_dir = Path(template_dir) if template_dir else Path(__file__).parent / "templates"
        self.template_dir.mkdir(exist_ok=True)
        
        # Create default templates if they don't exist
        self._create_default_templates()
        
        # Initialize visualizers
        self.model_viz = ModelVisualizer()
        self.weather_viz = WeatherVisualizer()
        
        # Setup Jinja2 environment
        self.env = Environment(loader=FileSystemLoader(str(self.template_dir)))
        
    def _create_default_templates(self):
        """Create default report templates"""
        # HTML template
        html_template = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/css/bootstrap.min.css" rel="stylesheet">
    <script src="https://cdn.plot.ly/plotly-latest.min.js"></script>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; }
        .metric-card { background: #f8f9fa; padding: 20px; margin: 10px; border-radius: 8px; }
        .chart-container { margin: 20px 0; }
        table { margin: 20px 0; }
        .summary-box { background: #e9ecef; padding: 15px; border-radius: 5px; margin: 20px 0; }
        .model-section { margin: 30px 0; padding: 20px; border: 1px solid #dee2e6; border-radius: 8px; }
    </style>
</head>
<body>
    <div class="container">
        <h1 class="text-center mb-4">{{ title }}</h1>
        <p class="text-center text-muted">Generated: {{ timestamp }}</p>
        
        <div class="summary-box">
            <h2>Executive Summary</h2>
            <p>{{ summary }}</p>
            <div class="row">
                <div class="col-md-3 metric-card">
                    <h4>Best Model</h4>
                    <h2 class="text-primary">{{ best_model }}</h2>
                </div>
                <div class="col-md-3 metric-card">
                    <h4>Best Accuracy</h4>
                    <h2 class="text-success">{{ "%.3f"|format(best_accuracy) }}</h2>
                </div>
                <div class="col-md-3 metric-card">
                    <h4>Models Trained</h4>
                    <h2 class="text-info">{{ n_models }}</h2>
                </div>
                <div class="col-md-3 metric-card">
                    <h4>Total Time</h4>
                    <h2 class="text-warning">{{ "%.1f"|format(total_time) }}h</h2>
                </div>
            </div>
        </div>
        
        <h2 class="mt-5">Model Comparison</h2>
        <div class="table-responsive">
            {{ comparison_table | safe }}
        </div>
        
        <h2 class="mt-5">Performance Visualizations</h2>
        {% for chart in charts %}
        <div class="chart-container">
            <h3>{{ chart.title }}</h3>
            {{ chart.html | safe }}
        </div>
        {% endfor %}
        
        <h2 class="mt-5">Detailed Model Results</h2>
        {% for model in models %}
        <div class="model-section">
            <h3>{{ model.name }}</h3>
            <div class="row">
                <div class="col-md-6">
                    <h5>Performance Metrics</h5>
                    <table class="table table-striped">
                        <tr><td>Accuracy</td><td>{{ "%.4f"|format(model.accuracy) }}</td></tr>
                        <tr><td>Precision</td><td>{{ "%.4f"|format(model.precision) }}</td></tr>
                        <tr><td>Recall</td><td>{{ "%.4f"|format(model.recall) }}</td></tr>
                        <tr><td>F1 Score</td><td>{{ "%.4f"|format(model.f1) }}</td></tr>
                        <tr><td>AUC</td><td>{{ "%.4f"|format(model.auc) if model.auc else "N/A" }}</td></tr>
                        <tr><td>Training Time</td><td>{{ "%.2f"|format(model.training_time) }}s</td></tr>
                    </table>
                </div>
                <div class="col-md-6">
                    {% if model.confusion_matrix %}
                    <h5>Confusion Matrix</h5>
                    <img src="{{ model.confusion_matrix }}" class="img-fluid">
                    {% endif %}
                </div>
            </div>
            {% if model.feature_importance %}
            <h5>Top Features</h5>
            <table class="table table-sm">
                <thead>
                    <tr><th>Feature</th><th>Importance</th></tr>
                </thead>
                <tbody>
                    {% for feature in model.feature_importance[:10] %}
                    <tr><td>{{ feature.name }}</td><td>{{ "%.4f"|format(feature.importance) }}</td></tr>
                    {% endfor %}
                </tbody>
            </table>
            {% endif %}
        </div>
        {% endfor %}
        
        <h2 class="mt-5">Conclusions and Recommendations</h2>
        <div class="summary-box">
            {{ conclusions | safe }}
        </div>
        
        <footer class="mt-5 text-center text-muted">
            <p>Weather Regulation Prediction System - Automated Report</p>
        </footer>
    </div>
</body>
</html>
"""
        
        # Markdown template
        markdown_template = """# {{ title }}

Generated: {{ timestamp }}

## Executive Summary

{{ summary }}

### Key Metrics
- **Best Model**: {{ best_model }}
- **Best Accuracy**: {{ "%.3f"|format(best_accuracy) }}
- **Models Trained**: {{ n_models }}
- **Total Training Time**: {{ "%.1f"|format(total_time) }} hours

## Model Comparison

{{ comparison_table }}

## Performance Analysis

{% for model in models %}
### {{ model.name }}

**Performance Metrics:**
- Accuracy: {{ "%.4f"|format(model.accuracy) }}
- Precision: {{ "%.4f"|format(model.precision) }}
- Recall: {{ "%.4f"|format(model.recall) }}
- F1 Score: {{ "%.4f"|format(model.f1) }}
- AUC: {{ "%.4f"|format(model.auc) if model.auc else "N/A" }}
- Training Time: {{ "%.2f"|format(model.training_time) }} seconds

{% if model.feature_importance %}
**Top 10 Features:**
{% for feature in model.feature_importance[:10] %}
1. {{ feature.name }}: {{ "%.4f"|format(feature.importance) }}
{% endfor %}
{% endif %}

{% endfor %}

## Conclusions and Recommendations

{{ conclusions }}

---
*Weather Regulation Prediction System - Automated Report*
"""
        
        # LaTeX template
        latex_template = r"""
\documentclass[11pt]{article}
\usepackage{graphicx}
\usepackage{booktabs}
\usepackage{hyperref}
\usepackage{float}
\usepackage{geometry}
\geometry{margin=1in}

\title{{ "{" }}{{ title }}{{ "}" }}
\author{Weather Regulation Prediction System}
\date{{ "{" }}{{ timestamp }}{{ "}" }}

\begin{document}
\maketitle

\section{Executive Summary}
{{ summary }}

\subsection{Key Metrics}
\begin{itemize}
    \item \textbf{Best Model}: {{ best_model }}
    \item \textbf{Best Accuracy}: {{ "%.3f"|format(best_accuracy) }}
    \item \textbf{Models Trained}: {{ n_models }}
    \item \textbf{Total Training Time}: {{ "%.1f"|format(total_time) }} hours
\end{itemize}

\section{Model Comparison}
{{ comparison_table_latex | safe }}

\section{Detailed Results}
{% for model in models %}
\subsection{{ "{" }}{{ model.name }}{{ "}" }}

\begin{table}[H]
\centering
\begin{tabular}{lr}
\toprule
Metric & Value \\
\midrule
Accuracy & {{ "%.4f"|format(model.accuracy) }} \\
Precision & {{ "%.4f"|format(model.precision) }} \\
Recall & {{ "%.4f"|format(model.recall) }} \\
F1 Score & {{ "%.4f"|format(model.f1) }} \\
AUC & {{ "%.4f"|format(model.auc) if model.auc else "N/A" }} \\
Training Time & {{ "%.2f"|format(model.training_time) }}s \\
\bottomrule
\end{tabular}
\caption{{ "{" }}Performance metrics for {{ model.name }}{{ "}" }}
\end{table}

{% endfor %}

\section{Conclusions and Recommendations}
{{ conclusions }}

\end{document}
"""
        
        # Save templates
        (self.template_dir / "report_html.jinja2").write_text(html_template)
        (self.template_dir / "report_markdown.jinja2").write_text(markdown_template)
        (self.template_dir / "report_latex.jinja2").write_text(latex_template)
    
    def generate_report(self, experiment: ExperimentResult,
                       format: str = 'html',
                       output_path: Optional[str] = None,
                       include_visualizations: bool = True) -> str:
        """
        Generate comprehensive report
        
        Args:
            experiment: Experiment results
            format: Output format ('html', 'pdf', 'markdown', 'latex', 'pptx')
            output_path: Output file path
            include_visualizations: Whether to include charts
            
        Returns:
            Path to generated report
        """
        # Prepare report data
        report_data = self._prepare_report_data(experiment, include_visualizations)
        
        # Generate report based on format
        if format == 'html':
            return self._generate_html_report(report_data, output_path)
        elif format == 'pdf':
            return self._generate_pdf_report(report_data, output_path)
        elif format == 'markdown':
            return self._generate_markdown_report(report_data, output_path)
        elif format == 'latex':
            return self._generate_latex_report(report_data, output_path)
        elif format == 'pptx':
            return self._generate_powerpoint_report(report_data, output_path)
        else:
            raise ValueError(f"Unsupported format: {format}")
    
    def _prepare_report_data(self, experiment: ExperimentResult, 
                           include_visualizations: bool) -> Dict[str, Any]:
        """Prepare data for report generation"""
        
        # Basic information
        data = {
            'title': f"Experiment Report: {experiment.experiment_name}",
            'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            'experiment_id': experiment.experiment_id,
            'best_model': experiment.best_model,
            'best_accuracy': experiment.best_accuracy or 0,
            'n_models': len(experiment.model_results),
            'total_time': experiment.total_training_time / 3600,  # Convert to hours
        }
        
        # Executive summary
        data['summary'] = self._generate_executive_summary(experiment)
        
        # Model comparison table
        if experiment.comparison_metrics is not None:
            data['comparison_table'] = experiment.comparison_metrics.to_html(
                classes='table table-striped', index=False
            )
            data['comparison_table_latex'] = experiment.comparison_metrics.to_latex(
                index=False, float_format="%.4f"
            )
        else:
            data['comparison_table'] = "No comparison data available"
            data['comparison_table_latex'] = "No comparison data available"
        
        # Detailed model results
        data['models'] = []
        for model_name, model_result in experiment.model_results.items():
            model_data = {
                'name': model_name,
                'accuracy': model_result.test_accuracy or 0,
                'precision': model_result.test_precision or 0,
                'recall': model_result.test_recall or 0,
                'f1': model_result.test_f1 or 0,
                'auc': model_result.test_auc,
                'training_time': model_result.training_time
            }
            
            # Feature importance
            if model_result.feature_importance is not None:
                model_data['feature_importance'] = [
                    {'name': row['feature'], 'importance': row['importance']}
                    for _, row in model_result.feature_importance.head(10).iterrows()
                ]
            
            # Confusion matrix (as base64 image for HTML)
            if include_visualizations and model_result.confusion_matrix is not None:
                cm_fig = self._create_confusion_matrix_figure(
                    model_result.confusion_matrix, model_name
                )
                model_data['confusion_matrix'] = self._fig_to_base64(cm_fig)
                plt.close(cm_fig)
            
            data['models'].append(model_data)
        
        # Visualizations
        if include_visualizations:
            data['charts'] = self._create_visualizations(experiment)
        
        # Conclusions
        data['conclusions'] = self._generate_conclusions(experiment)
        
        return data
    
    def _generate_executive_summary(self, experiment: ExperimentResult) -> str:
        """Generate executive summary"""
        summary = f"""
        This report presents the results of the weather regulation prediction experiment 
        '{experiment.experiment_name}' conducted on {experiment.timestamp.strftime('%Y-%m-%d')}.
        A total of {len(experiment.model_results)} models were trained and evaluated.
        The best performing model was {experiment.best_model} with an accuracy of 
        {experiment.best_accuracy:.3f}. The experiment took {experiment.total_training_time/3600:.1f} 
        hours to complete.
        """
        return summary.strip()
    
    def _generate_conclusions(self, experiment: ExperimentResult) -> str:
        """Generate conclusions and recommendations"""
        conclusions = []
        
        # Best model analysis
        conclusions.append(f"### Best Model: {experiment.best_model}")
        conclusions.append(f"The {experiment.best_model} model achieved the highest accuracy "
                         f"of {experiment.best_accuracy:.3f}.")
        
        # Performance insights
        if experiment.comparison_metrics is not None:
            avg_f1 = experiment.comparison_metrics['f1_score'].mean()
            conclusions.append(f"\nThe average F1 score across all models was {avg_f1:.3f}, "
                             f"indicating {'good' if avg_f1 > 0.8 else 'moderate'} overall performance.")
        
        # Recommendations
        conclusions.append("\n### Recommendations")
        
        if experiment.best_accuracy < 0.9:
            conclusions.append("- Consider collecting more training data to improve model performance")
            conclusions.append("- Experiment with advanced feature engineering techniques")
            conclusions.append("- Try ensemble methods to combine multiple models")
        else:
            conclusions.append("- The model shows excellent performance and is ready for deployment")
            conclusions.append("- Monitor model performance over time to detect any degradation")
            conclusions.append("- Consider A/B testing in production environment")
        
        # Efficiency analysis
        if experiment.total_training_time > 7200:  # 2 hours
            conclusions.append("\n- Training time was significant. Consider:")
            conclusions.append("  - Using distributed training for faster iterations")
            conclusions.append("  - Implementing early stopping to reduce training time")
            conclusions.append("  - Optimizing hyperparameter search space")
        
        return "\n".join(conclusions)
    
    def _create_confusion_matrix_figure(self, cm: np.ndarray, title: str) -> plt.Figure:
        """Create confusion matrix figure"""
        fig, ax = plt.subplots(figsize=(6, 5))
        sns.heatmap(cm, annot=True, fmt='d', cmap='Blues', ax=ax)
        ax.set_title(f'Confusion Matrix - {title}')
        ax.set_xlabel('Predicted')
        ax.set_ylabel('Actual')
        plt.tight_layout()
        return fig
    
    def _fig_to_base64(self, fig: plt.Figure) -> str:
        """Convert matplotlib figure to base64 string"""
        buffer = BytesIO()
        fig.savefig(buffer, format='png', dpi=100, bbox_inches='tight')
        buffer.seek(0)
        img_str = base64.b64encode(buffer.read()).decode()
        return f"data:image/png;base64,{img_str}"
    
    def _create_visualizations(self, experiment: ExperimentResult) -> List[Dict[str, str]]:
        """Create visualization charts"""
        charts = []
        
        # Model comparison bar chart
        if experiment.comparison_metrics is not None:
            import plotly.express as px
            
            # Accuracy comparison
            fig = px.bar(
                experiment.comparison_metrics,
                x='model',
                y='accuracy',
                title='Model Accuracy Comparison',
                height=400
            )
            charts.append({
                'title': 'Model Accuracy Comparison',
                'html': fig.to_html(include_plotlyjs=False, div_id="accuracy-chart")
            })
            
            # Multi-metric comparison
            metrics = ['precision', 'recall', 'f1_score']
            if all(m in experiment.comparison_metrics.columns for m in metrics):
                fig = px.bar(
                    experiment.comparison_metrics.melt(
                        id_vars=['model'],
                        value_vars=metrics,
                        var_name='metric',
                        value_name='score'
                    ),
                    x='model',
                    y='score',
                    color='metric',
                    barmode='group',
                    title='Multi-Metric Comparison',
                    height=400
                )
                charts.append({
                    'title': 'Multi-Metric Comparison',
                    'html': fig.to_html(include_plotlyjs=False, div_id="metrics-chart")
                })
        
        return charts
    
    def _generate_html_report(self, data: Dict[str, Any], output_path: Optional[str]) -> str:
        """Generate HTML report"""
        template = self.env.get_template("report_html.jinja2")
        html_content = template.render(**data)
        
        if output_path is None:
            output_path = f"report_{data['experiment_id']}.html"
        
        Path(output_path).write_text(html_content)
        return output_path
    
    def _generate_markdown_report(self, data: Dict[str, Any], output_path: Optional[str]) -> str:
        """Generate Markdown report"""
        template = self.env.get_template("report_markdown.jinja2")
        md_content = template.render(**data)
        
        if output_path is None:
            output_path = f"report_{data['experiment_id']}.md"
        
        Path(output_path).write_text(md_content)
        return output_path
    
    def _generate_latex_report(self, data: Dict[str, Any], output_path: Optional[str]) -> str:
        """Generate LaTeX report"""
        template = self.env.get_template("report_latex.jinja2")
        latex_content = template.render(**data)
        
        if output_path is None:
            output_path = f"report_{data['experiment_id']}.tex"
        
        Path(output_path).write_text(latex_content)
        return output_path
    
    def _generate_pdf_report(self, data: Dict[str, Any], output_path: Optional[str]) -> str:
        """Generate PDF report using ReportLab"""
        if not REPORTLAB_AVAILABLE:
            raise ImportError("ReportLab is required for PDF generation. Install with: pip install reportlab")
        
        if output_path is None:
            output_path = f"report_{data['experiment_id']}.pdf"
        
        # Create PDF document
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1f77b4'),
            spaceAfter=30,
            alignment=1  # Center
        )
        story.append(Paragraph(data['title'], title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Timestamp
        story.append(Paragraph(f"Generated: {data['timestamp']}", styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", styles['Heading2']))
        story.append(Paragraph(data['summary'], styles['Normal']))
        story.append(Spacer(1, 0.2*inch))
        
        # Key Metrics
        metrics_data = [
            ['Metric', 'Value'],
            ['Best Model', data['best_model']],
            ['Best Accuracy', f"{data['best_accuracy']:.3f}"],
            ['Models Trained', str(data['n_models'])],
            ['Total Time', f"{data['total_time']:.1f} hours"]
        ]
        
        metrics_table = Table(metrics_data, colWidths=[3*inch, 2*inch])
        metrics_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 14),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(metrics_table)
        story.append(PageBreak())
        
        # Model Results
        story.append(Paragraph("Model Results", styles['Heading2']))
        for model in data['models']:
            story.append(Paragraph(model['name'], styles['Heading3']))
            
            model_data = [
                ['Metric', 'Value'],
                ['Accuracy', f"{model['accuracy']:.4f}"],
                ['Precision', f"{model['precision']:.4f}"],
                ['Recall', f"{model['recall']:.4f}"],
                ['F1 Score', f"{model['f1']:.4f}"],
                ['Training Time', f"{model['training_time']:.2f}s"]
            ]
            
            model_table = Table(model_data, colWidths=[2*inch, 1.5*inch])
            model_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.lightblue),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey)
            ]))
            story.append(model_table)
            story.append(Spacer(1, 0.2*inch))
        
        # Conclusions
        story.append(PageBreak())
        story.append(Paragraph("Conclusions and Recommendations", styles['Heading2']))
        conclusions_text = data['conclusions'].replace('\n', '<br/>')
        story.append(Paragraph(conclusions_text, styles['Normal']))
        
        # Build PDF
        doc.build(story)
        return output_path
    
    def _generate_powerpoint_report(self, data: Dict[str, Any], output_path: Optional[str]) -> str:
        """Generate PowerPoint presentation"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PowerPoint generation. Install with: pip install python-pptx")
        
        if output_path is None:
            output_path = f"report_{data['experiment_id']}.pptx"
        
        # Create presentation
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = data['title']
        subtitle.text = f"Generated: {data['timestamp']}"
        
        # Executive Summary slide
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        summary_slide.shapes.title.text = "Executive Summary"
        summary_slide.placeholders[1].text = data['summary']
        
        # Key Metrics slide
        metrics_slide = prs.slides.add_slide(prs.slide_layouts[5])
        metrics_slide.shapes.title.text = "Key Metrics"
        
        # Add metrics as bullet points
        body_shape = metrics_slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = f"Best Model: {data['best_model']}"
        
        metrics = [
            f"Best Accuracy: {data['best_accuracy']:.3f}",
            f"Models Trained: {data['n_models']}",
            f"Total Training Time: {data['total_time']:.1f} hours"
        ]
        
        for metric in metrics:
            p = tf.add_paragraph()
            p.text = metric
            p.level = 0
        
        # Model Results slides
        for model in data['models']:
            model_slide = prs.slides.add_slide(prs.slide_layouts[1])
            model_slide.shapes.title.text = f"Model: {model['name']}"
            
            body = model_slide.placeholders[1]
            tf = body.text_frame
            tf.text = "Performance Metrics:"
            
            metrics = [
                f"Accuracy: {model['accuracy']:.4f}",
                f"Precision: {model['precision']:.4f}",
                f"Recall: {model['recall']:.4f}",
                f"F1 Score: {model['f1']:.4f}",
                f"Training Time: {model['training_time']:.2f}s"
            ]
            
            for metric in metrics:
                p = tf.add_paragraph()
                p.text = metric
                p.level = 1
        
        # Conclusions slide
        conclusions_slide = prs.slides.add_slide(prs.slide_layouts[1])
        conclusions_slide.shapes.title.text = "Conclusions & Recommendations"
        
        # Split conclusions into bullet points
        conclusions_lines = data['conclusions'].split('\n')
        body = conclusions_slide.placeholders[1]
        tf = body.text_frame
        tf.text = conclusions_lines[0] if conclusions_lines else ""
        
        for line in conclusions_lines[1:]:
            if line.strip():
                p = tf.add_paragraph()
                p.text = line.strip('- ')
                p.level = 0 if line.startswith('#') else 1
        
        # Save presentation
        prs.save(output_path)
        return output_path


# Utility functions
def generate_experiment_report(experiment: ExperimentResult,
                             format: str = 'html',
                             output_dir: str = "./reports") -> str:
    """Quick function to generate experiment report"""
    generator = ReportGenerator()
    
    output_path = Path(output_dir) / f"report_{experiment.experiment_id}.{format}"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    return generator.generate_report(experiment, format, str(output_path))


def generate_comparison_report(experiments: List[ExperimentResult],
                             output_path: str = "./reports/comparison_report.html") -> str:
    """Generate report comparing multiple experiments"""
    # This would create a specialized comparison report
    # Implementation would be similar but focus on comparing experiments
    pass